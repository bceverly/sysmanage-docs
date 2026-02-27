#!/usr/bin/env python3
"""Generate SysManage Overview Presentation (PPTX)."""

import os
import cairosvg
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ──────────────────────────────────────────────────────────────
DOCS_DIR = os.path.expanduser("~/dev/sysmanage-docs")
ASSETS_DIR = os.path.join(DOCS_DIR, "assets", "images")
OUT_DIR = os.path.join(DOCS_DIR, "Presentations")
OUT_FILE = os.path.join(OUT_DIR, "SysManage.pptx")

LOGO_SVG = os.path.join(ASSETS_DIR, "sysmanage-logo.svg")
ICON_SVG = os.path.join(ASSETS_DIR, "sysmanage-icon.svg")
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
LOGO_PNG = os.path.join(SCRIPT_DIR, "sysmanage-logo.png")
ICON_PNG = os.path.join(SCRIPT_DIR, "sysmanage-icon.png")
BRYAN_PHOTO = os.path.join(SCRIPT_DIR, "bryan-everly.jpg")
FEDOR_PHOTO = os.path.join(SCRIPT_DIR, "fedor-dikarev.jpg")

# ── Branding colors ───────────────────────────────────────────────────
PRIMARY_BLUE = RGBColor(0x19, 0x76, 0xD2)
DARK_BLUE = RGBColor(0x0D, 0x3D, 0x5C)
GREEN = RGBColor(0x38, 0x8E, 0x3C)
BUTTON_GREEN = RGBColor(0x2E, 0x7D, 0x47)
DARK_TEXT = RGBColor(0x1A, 0x23, 0x32)
LIGHT_BG = RGBColor(0xF5, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BLUE = RGBColor(0xBB, 0xDE, 0xFB)
VERY_LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)

FONT = "Arial"

# ── Slide dimensions (16:9) ───────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────────────

def convert_svgs():
    """Convert SVG logos to PNG for embedding."""
    cairosvg.svg2png(url=LOGO_SVG, write_to=LOGO_PNG, output_width=660, output_height=180)
    cairosvg.svg2png(url=ICON_SVG, write_to=ICON_PNG, output_width=360, output_height=360)


def add_dark_bg(slide):
    """Fill the entire slide with DARK_BLUE background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE


def add_light_bg(slide):
    """Fill the slide with light background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG


def add_bottom_bar(slide):
    """Add a thin branded bar at the bottom of the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), SLIDE_H - Inches(0.35),
        SLIDE_W, Inches(0.35),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()


def add_top_bar(slide):
    """Add a branded header bar at top."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_W, Inches(1.1),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()


def add_footer_text(slide, text="SysManage  |  sysmanage.org"):
    """Add small footer text in the bottom bar."""
    txBox = slide.shapes.add_textbox(
        Inches(0.5), SLIDE_H - Inches(0.33),
        Inches(6), Inches(0.3),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.name = FONT


def set_run(paragraph, text, size=18, bold=False, color=DARK_TEXT, font=FONT):
    """Add a run with formatting to a paragraph."""
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return run


def add_title_textbox(slide, text, left=0.8, top=0.15, width=11, height=0.8,
                      size=36, color=WHITE, bold=True):
    """Add a title text in the header bar area."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_run(p, text, size=size, bold=bold, color=color)
    return tf


def add_body_textbox(slide, left=0.8, top=1.4, width=11.5, height=5.2):
    """Add a textbox for body content."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_bullet_slide(slide, title, bullets, sub_bullets=None):
    """Standard content slide with header bar, title, bullets, footer."""
    add_light_bg(slide)
    add_top_bar(slide)
    add_title_textbox(slide, title)
    add_bottom_bar(slide)
    add_footer_text(slide)

    tf = add_body_textbox(slide)
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(4)
        p.level = 0
        set_run(p, f"\u2022  {bullet}", size=18, color=DARK_TEXT)

        # Add sub-bullets if provided for this index
        if sub_bullets and i in sub_bullets:
            for sb in sub_bullets[i]:
                sp = tf.add_paragraph()
                sp.space_before = Pt(2)
                sp.space_after = Pt(2)
                sp.level = 1
                set_run(sp, f"    \u2013  {sb}", size=15, color=MEDIUM_GRAY)


def add_table_slide(slide, title, headers, rows, col_widths=None):
    """Content slide with a styled table."""
    add_light_bg(slide)
    add_top_bar(slide)
    add_title_textbox(slide, title)
    add_bottom_bar(slide)
    add_footer_text(slide)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(11.5)
    height = Inches(0.4 * n_rows)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = FONT

    # Data rows
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else VERY_LIGHT_BLUE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = DARK_TEXT
                paragraph.font.name = FONT


def add_accent_rect(slide, left, top, width, height, color=PRIMARY_BLUE):
    """Add a colored rectangle accent."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_in_shape(shape, text, size=14, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER):
    """Put text inside a shape."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT


def add_multiline_in_shape(shape, lines, size=12, color=WHITE, bold=False,
                           alignment=PP_ALIGN.LEFT):
    """Put multiple lines of text inside a shape."""
    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = FONT


# ── Slide Builders ─────────────────────────────────────────────────────

def slide_01_title(prs):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_dark_bg(slide)

    # Logo
    slide.shapes.add_picture(LOGO_PNG, Inches(4.0), Inches(1.0), Inches(5.3))

    # Main title
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(10), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, "Cross-Platform System Management", size=32, bold=True, color=WHITE)

    # Tagline
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    set_run(p2, "Manage every host, every OS, from one platform.", size=20, color=LIGHT_BLUE)

    # Date
    txBox3 = slide.shapes.add_textbox(Inches(3), Inches(5.3), Inches(7), Inches(0.6))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    set_run(p3, "February 2026", size=16, color=MEDIUM_GRAY)

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.0), Inches(5.0), Inches(3.3), Inches(0.04),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_BLUE
    line.line.fill.background()


def slide_02_agenda(prs):
    """Agenda slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_top_bar(slide)
    add_title_textbox(slide, "Agenda")
    add_bottom_bar(slide)
    add_footer_text(slide)

    items = [
        ("Overview", "What is SysManage & architecture", "~10 min"),
        ("Features", "Open source + Pro+ capabilities", "~8 min"),
        ("Security", "Security-first design & dev process", "~5 min"),
        ("Live Demo", "Hands-on walkthrough", "~30 min"),
        ("Futures", "Roadmap to v3.0.0.0", "~5 min"),
        ("Recap", "Key takeaways & resources", "~2 min"),
    ]

    for i, (label, desc, time) in enumerate(items):
        y = 1.5 + i * 0.9
        # Number circle
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.2), Inches(y), Inches(0.55), Inches(0.55),
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = PRIMARY_BLUE
        circ.line.fill.background()
        add_text_in_shape(circ, str(i + 1), size=16, color=WHITE, bold=True)

        # Label
        txBox = slide.shapes.add_textbox(Inches(2.0), Inches(y), Inches(3), Inches(0.35))
        tf = txBox.text_frame
        set_run(tf.paragraphs[0], label, size=20, bold=True, color=DARK_BLUE)

        # Description
        txBox2 = slide.shapes.add_textbox(Inches(2.0), Inches(y + 0.32), Inches(5), Inches(0.3))
        tf2 = txBox2.text_frame
        set_run(tf2.paragraphs[0], desc, size=14, color=MEDIUM_GRAY)

        # Time badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(10.0), Inches(y + 0.05), Inches(1.3), Inches(0.4),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = VERY_LIGHT_BLUE
        badge.line.fill.background()
        add_text_in_shape(badge, time, size=12, color=DARK_BLUE, bold=True)


def slide_03_problem(prs):
    """The Problem slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "Heterogeneous OS estates \u2014 Linux, macOS, Windows, *BSD all in one environment",
        "Security & compliance requirements \u2014 SOC2, HIPAA, FedRAMP, CIS, DISA STIG",
        "Agent sprawl \u2014 separate tools for monitoring, patching, compliance, secrets",
        "Manual toil \u2014 SSH loops, ad-hoc scripts, no audit trail",
        "No single pane of glass \u2014 context-switching between dashboards",
        "Commercial lock-in \u2014 expensive per-host licensing with opaque code",
    ]
    add_bullet_slide(slide, "The Problem", bullets)


def slide_04_what_is(prs):
    """What is SysManage? slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_top_bar(slide)
    add_title_textbox(slide, "What is SysManage?")
    add_bottom_bar(slide)
    add_footer_text(slide)

    # Left column: description
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(6.5), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    items = [
        ("Server + Agent architecture", "Central server manages distributed agents on every host"),
        ("Real-time WebSocket communication", "Live metrics, status updates, and command execution"),
        ("FastAPI backend + React frontend", "Modern Python API with responsive web UI"),
        ("PostgreSQL database", "Reliable, scalable data store for all host & config data"),
        ("AGPLv3 open source core", "Transparent, auditable, community-driven"),
        ("Pro+ commercial modules", "Enterprise features with Cython-compiled add-ons"),
    ]

    for i, (title, desc) in enumerate(items):
        if i > 0:
            p = tf.add_paragraph()
            p.space_before = Pt(10)
        else:
            p = tf.paragraphs[0]
            p.space_before = Pt(4)
        set_run(p, f"\u2022  {title}", size=17, bold=True, color=DARK_BLUE)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        set_run(p2, f"     {desc}", size=14, color=MEDIUM_GRAY)

    # Right column: icon
    slide.shapes.add_picture(ICON_PNG, Inches(9.0), Inches(2.0), Inches(3.0))


def slide_05_who_bryan(prs):
    """Who is SysManage? — Bryan Everly slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_top_bar(slide)
    add_title_textbox(slide, "Who is SysManage? \u2014 Bryan Everly")
    add_bottom_bar(slide)
    add_footer_text(slide)

    # Photo
    slide.shapes.add_picture(BRYAN_PHOTO, Inches(10.5), Inches(1.25), Inches(1.5))

    # Role badge
    badge = add_accent_rect(slide, 0.8, 1.35, 3.0, 0.5, PRIMARY_BLUE)
    add_text_in_shape(badge, "Founder & CEO", size=16, color=WHITE, bold=True)

    # Career timeline
    career = [
        ("Software Artistry", "Early employee \u2192 Dir. of Worldwide Product Dev  \u2022  IPO 1995  \u2022  Sold to IBM/Tivoli 1998"),
        ("IBM", "Led global engineering team post-acquisition (3 years)"),
        ("SaaS Founder", "Founded & sold a SaaS company securing sensitive HR data (~8 years)"),
        ("ExactTarget", "Engineering leader (1/3 of technology team)  \u2022  Acquired by Salesforce (now Salesforce Marketing Cloud)"),
        ("Aprimo / Teradata", "VP of Worldwide Engineering  \u2022  $525M acquisition by Teradata"),
        ("NextGear Capital", "CTO  \u2022  Built dev org from scratch  \u2022  Grew loan volume $1.6B \u2192 $3.2B  \u2022  CTO of the Year (IBJ/TechPoint)"),
        ("Cox Automotive", "CTO & CISO"),
        ("Cummins", "Chief Enterprise Architect & Exec. Dir. Software Engineering"),
        ("Canonical", "VP of Software Engineering  \u2022  Ubuntu, IoT, real-time Linux, HPC"),
    ]

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (company, detail) in enumerate(career):
        if i > 0:
            p = tf.add_paragraph()
            p.space_before = Pt(6)
        else:
            p = tf.paragraphs[0]
            p.space_before = Pt(2)
        set_run(p, f"{company}  ", size=15, bold=True, color=DARK_BLUE)
        set_run(p, f"\u2014  {detail}", size=13, bold=False, color=MEDIUM_GRAY)

    # Education & OSS footer
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    set_run(tf2.paragraphs[0],
            "M.S. Cybersecurity, UC Berkeley  \u2022  B.S. Computer Science, Indiana State  \u2022  OpenBSD Ports Maintainer",
            size=12, color=DARK_TEXT)


def slide_06_who_fedor(prs):
    """Who is SysManage? — Fedor Dikarev slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_top_bar(slide)
    add_title_textbox(slide, "Who is SysManage? \u2014 Fedor Dikarev")
    add_bottom_bar(slide)
    add_footer_text(slide)

    # Photo
    slide.shapes.add_picture(FEDOR_PHOTO, Inches(10.5), Inches(1.25), Inches(1.5))

    # Role badge
    badge = add_accent_rect(slide, 0.8, 1.35, 4.5, 0.5, PRIMARY_BLUE)
    add_text_in_shape(badge, "Lead SRE / CI/CD Engineer", size=16, color=WHITE, bold=True)

    # Background items
    items = [
        ("NGINX",
         "Infrastructure and operations engineering at NGINX, the world's most widely deployed web server and reverse proxy"),
        ("Databricks (via Neon acquisition)",
         "Infrastructure engineer at Databricks following the ~$1B acquisition of Neon, the serverless Postgres company"),
        ("Neon \u2014 Serverless Postgres",
         "Member of Technical Staff  \u2022  Led CI/CD infrastructure migration and cost optimization  \u2022  Reduced cloud CI costs by 50%"),
        ("Infrastructure & DevOps Expertise",
         "Deep experience with AWS, Docker, Kubernetes, Prometheus monitoring, and high-availability architectures"),
        ("CI/CD & Automation",
         "Designed and operated large-scale CI pipelines  \u2022  GitHub Actions, ephemeral runners, caching optimization"),
        ("Based in Germany",
         "European engineering perspective and timezone coverage"),
    ]

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.5), Inches(4.2))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (title, desc) in enumerate(items):
        if i > 0:
            p = tf.add_paragraph()
            p.space_before = Pt(10)
        else:
            p = tf.paragraphs[0]
            p.space_before = Pt(4)
        set_run(p, f"\u2022  {title}", size=16, bold=True, color=DARK_BLUE)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        set_run(p2, f"     {desc}", size=13, color=MEDIUM_GRAY)


def slide_07_growth_strategy(prs):
    """Growth Strategy — VC & Advisors slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_top_bar(slide)
    add_title_textbox(slide, "Growth Strategy")
    add_bottom_bar(slide)
    add_footer_text(slide)

    # VC box
    vc_box = add_accent_rect(slide, 0.8, 1.5, 5.5, 3.0, DARK_BLUE)
    add_multiline_in_shape(vc_box, [
        "Venture Capital",
        "",
        "In active discussions with a European",
        "venture capital fund focused on open",
        "source and infrastructure software",
    ], size=14, color=WHITE, alignment=PP_ALIGN.LEFT)

    # Advisors box
    adv_box = add_accent_rect(slide, 7.0, 1.5, 5.5, 3.0, PRIMARY_BLUE)
    add_multiline_in_shape(adv_box, [
        "Advisory Network",
        "",
        "Engaging potential advisors from",
        "proven open source \u2192 commercial",
        "success stories including NGINX",
        "and WordPress ecosystems",
    ], size=14, color=WHITE, alignment=PP_ALIGN.LEFT)

    # Bottom: strategy note
    strategy_box = add_accent_rect(slide, 0.8, 5.0, 11.7, 1.5, BUTTON_GREEN)
    add_multiline_in_shape(strategy_box, [
        "Open Source \u2192 Commercial Playbook",
        "",
        "AGPLv3 core builds community trust & adoption  \u2192  Pro+ tiers monetize enterprise needs",
        "Following the proven model: open core + commercial extensions + enterprise support",
    ], size=13, color=WHITE, alignment=PP_ALIGN.CENTER)


def slide_08_architecture(prs):
    """Architecture Diagram slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_top_bar(slide)
    add_title_textbox(slide, "Architecture")
    add_bottom_bar(slide)
    add_footer_text(slide)

    # Agent box (left)
    agent_box = add_accent_rect(slide, 0.8, 2.0, 3.2, 2.0, GREEN)
    add_multiline_in_shape(agent_box, [
        "SysManage Agent",
        "",
        "Local DB (SQLite)",
        "System Metrics",
    ], size=12, color=WHITE, alignment=PP_ALIGN.CENTER)

    # mTLS label (centered between agent and server)
    mtls_box = slide.shapes.add_textbox(Inches(4.1), Inches(2.4), Inches(1.6), Inches(0.8))
    tf = mtls_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, "mTLS", size=13, bold=True, color=PRIMARY_BLUE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    set_run(p2, "\u2190 \u2192", size=16, bold=True, color=PRIMARY_BLUE)

    # Server box
    server_box = add_accent_rect(slide, 5.8, 1.5, 4.3, 3.0, DARK_BLUE)
    add_multiline_in_shape(server_box, [
        "SysManage Server",
        "",
        "FastAPI Backend",
        "React Frontend",
        "WebSocket Engine",
    ], size=13, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER)

    # PostgreSQL
    db_box = add_accent_rect(slide, 5.8, 5.0, 4.3, 0.8, PRIMARY_BLUE)
    add_text_in_shape(db_box, "PostgreSQL Database", size=14, color=WHITE, bold=True)

    # Pro+ modules
    pro_box = add_accent_rect(slide, 10.5, 1.5, 2.3, 3.0, BUTTON_GREEN)
    add_multiline_in_shape(pro_box, [
        "Pro+",
        "Modules",
        "",
        "Health",
        "Compliance",
        "Reporting",
        "Secrets",
        "CVE",
    ], size=10, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Message queuing note
    queue_box = add_accent_rect(slide, 0.8, 4.6, 4.6, 1.2, DARK_BLUE)
    add_multiline_in_shape(queue_box, [
        "Message Queue Architecture",
        "All agent \u2194 server communication is queued",
        "to the database. Background threads process",
        "messages \u2014 no data loss on crash or restart.",
    ], size=10, color=WHITE, alignment=PP_ALIGN.CENTER)


def slide_09_cross_platform(prs):
    """Cross-Platform Matrix slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    headers = ["Operating System", "Server", "Agent"]
    rows = [
        ("Linux (Ubuntu, Debian, RHEL, Alpine, etc.)", "\u2713", "\u2713"),
        ("macOS (13+)", "\u2713", "\u2713"),
        ("Windows (10/11, Server 2019+)", "\u2713", "\u2713"),
        ("FreeBSD (13+)", "\u2713", "\u2713"),
        ("OpenBSD (7.4+)", "\u2713", "\u2713"),
        ("NetBSD (10+)", "\u2713", "\u2713"),
    ]
    add_table_slide(slide, "Cross-Platform Support", headers, rows,
                    col_widths=[6.0, 2.5, 2.5])


def slide_10_packaging(prs):
    """Packaging Formats & OS Versions slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    headers = ["Format", "OS / Distro", "Versions"]
    rows = [
        (".deb", "Ubuntu, Debian", "20.04+, Bookworm+"),
        (".rpm", "RHEL, Fedora, SUSE", "8+, 38+, 15+"),
        (".pkg", "macOS", "13 Ventura+"),
        (".msi", "Windows", "10/11, Server 2019+"),
        (".snap", "Ubuntu/Snapcraft", "Core 22+"),
        (".flatpak", "Linux (universal)", "Flathub"),
        (".apk", "Alpine Linux", "3.18+"),
        (".tgz", "FreeBSD", "13+"),
        ("port", "OpenBSD, NetBSD", "7.4+ / 10+"),
    ]
    add_table_slide(slide, "Packaging Formats & OS Versions", headers, rows,
                    col_widths=[2.0, 4.0, 4.0])


def slide_11_i18n(prs):
    """Internationalization slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_top_bar(slide)
    add_title_textbox(slide, "Internationalization")
    add_bottom_bar(slide)
    add_footer_text(slide)

    languages = [
        ("English", "en"), ("French", "fr"), ("German", "de"),
        ("Spanish", "es"), ("Italian", "it"), ("Portuguese", "pt"),
        ("Dutch", "nl"), ("Russian", "ru"), ("Chinese (Simplified)", "zh"),
        ("Japanese", "ja"), ("Korean", "ko"), ("Hindi", "hi"),
        ("Arabic (RTL)", "ar"), ("Turkish", "tr"),
    ]

    cols = 3
    for i, (lang, code) in enumerate(languages):
        col = i % cols
        row = i // cols
        x = 1.0 + col * 4.0
        y = 1.6 + row * 0.8

        badge = add_accent_rect(slide, x, y, 3.4, 0.55, PRIMARY_BLUE if code != "ar" else GREEN)
        add_text_in_shape(badge, f"{lang}  ({code})", size=14, color=WHITE, bold=True)

    # Note
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(10), Inches(0.5))
    tf = txBox.text_frame
    set_run(tf.paragraphs[0], "Full UI + backend message localization  \u2022  RTL support for Arabic",
            size=14, color=MEDIUM_GRAY)


def slide_12_oss_features(prs):
    """Open Source Features slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "Host management \u2014 register, organize, tag, and monitor all hosts",
        "Software inventory \u2014 full package listing per host",
        "Update detection \u2014 identify available OS and package updates",
        "Package management \u2014 install, remove, upgrade packages remotely",
        "Repository management \u2014 configure package sources across hosts",
        "Certificate monitoring \u2014 track TLS certificate expiration dates",
        "Firewall & AV status \u2014 verify security posture at a glance",
        "Role-Based Access Control \u2014 Admin, Operator, Viewer roles",
        "User management \u2014 create, modify, deactivate user accounts",
        "Tagging system \u2014 flexible host categorization and filtering",
    ]
    add_bullet_slide(slide, "Open Source Features", bullets)


def slide_13_oss_monitoring(prs):
    """Open Source: Monitoring & Management slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "Real-time CPU, RAM, disk, and network metrics via WebSocket",
        "System uptime and last-seen tracking",
        "OS identification and version detection",
        "Storage inventory \u2014 partitions, mount points, usage",
        "Ubuntu Pro integration \u2014 ESM status and entitlements",
        "Host detail dashboard with at-a-glance health overview",
        "WebSocket-driven live updates \u2014 no polling, no page refresh",
    ]
    add_bullet_slide(slide, "Open Source: Monitoring & Management", bullets)


def slide_14_proplus_licensing(prs):
    """Pro+ Licensing Model slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_top_bar(slide)
    add_title_textbox(slide, "Pro+ Licensing Model")
    add_bottom_bar(slide)
    add_footer_text(slide)

    # Professional tier box
    pro_box = add_accent_rect(slide, 0.8, 1.6, 5.5, 4.5, PRIMARY_BLUE)
    add_multiline_in_shape(pro_box, [
        "Professional Tier",
        "",
        "\u2022  Health Engine (AI-powered)",
        "\u2022  Compliance Engine (CIS/STIG)",
        "\u2022  Reporting Engine (PDF/HTML)",
        "\u2022  Audit Engine (tamper-evident)",
        "\u2022  Secrets Engine (OpenBAO)",
        "\u2022  Container Engine (LXD/WSL)",
        "\u2022  CVE Vulnerability Engine",
        "\u2022  Alerting Engine",
    ], size=14, color=WHITE, alignment=PP_ALIGN.LEFT)

    # Enterprise tier box
    ent_box = add_accent_rect(slide, 7.0, 1.6, 5.5, 4.5, DARK_BLUE)
    add_multiline_in_shape(ent_box, [
        "Enterprise Tier",
        "",
        "Everything in Professional, plus:",
        "",
        "\u2022  AV Management",
        "\u2022  Firewall Orchestration",
        "\u2022  Virtualization (KVM/bhyve/VMM)",
        "\u2022  Observability",
        "\u2022  Automation & Fleet",
        "\u2022  Air-Gapped Support",
        "\u2022  MFA",
    ], size=14, color=WHITE, alignment=PP_ALIGN.LEFT)

    # Bottom note
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11), Inches(0.5))
    tf = txBox.text_frame
    set_run(tf.paragraphs[0],
            "ECDSA P-521 license validation  \u2022  Feature gating  \u2022  Host count enforcement  \u2022  Cython-compiled modules",
            size=13, color=MEDIUM_GRAY)


def slide_15_proplus_professional(prs):
    """Pro+ Professional Tier slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "proplus_core \u2014 license validation, feature gating, module loader",
        "health_engine \u2014 AI-powered system health analysis and recommendations",
        "compliance_engine \u2014 CIS Benchmarks and DISA STIG auditing",
        "reporting_engine \u2014 PDF and HTML report generation",
        "audit_engine \u2014 tamper-evident audit logging with hash chains",
        "secrets_engine \u2014 OpenBAO-backed encrypted secrets management",
        "container_engine \u2014 LXD and WSL container management",
    ]
    add_bullet_slide(slide, "Pro+ Professional Tier", bullets)


def slide_16_proplus_enterprise(prs):
    """Pro+ Enterprise Tier slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "vuln_engine \u2014 CVE vulnerability scanning and tracking",
        "alerting_engine \u2014 Email, Webhook, Slack, and Teams notifications",
    ]
    sub = {
        2: [
            "AV management \u2014 centralized antivirus policy and status",
            "Firewall orchestration \u2014 cross-platform firewall rule management",
            "Virtualization \u2014 KVM, bhyve, VMM hypervisor management",
            "Observability \u2014 centralized log and metric aggregation",
            "Automation & Fleet \u2014 task scheduling and fleet-wide operations",
            "Air-gapped support \u2014 offline/disconnected environment management",
            "MFA \u2014 multi-factor authentication for the management console",
        ],
    }
    add_bullet_slide(slide, "Pro+ Enterprise Tier", bullets + ["Upcoming enterprise modules:"], sub)


def slide_17_ecosystem(prs):
    """Ecosystem Integrations slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_top_bar(slide)
    add_title_textbox(slide, "Ecosystem Integrations")
    add_bottom_bar(slide)
    add_footer_text(slide)

    # Categories with integrations — two-column layout to fit vertically
    left_categories = [
        ("Secrets & Key Mgmt", [
            ("OpenBAO / Vault", "Encrypted secrets, API keys"),
        ], DARK_BLUE),
        ("Observability", [
            ("OpenTelemetry", "Distributed tracing & metrics"),
            ("Grafana", "Dashboard visualization"),
            ("Prometheus", "Metrics collection & querying"),
            ("Graylog", "Log aggregation (GELF, syslog)"),
        ], PRIMARY_BLUE),
        ("Vulnerability Data", [
            ("NIST NVD", "National Vulnerability Database"),
            ("Ubuntu / Debian / RHEL / MSRC / FreeBSD", "OS-specific advisories"),
        ], BUTTON_GREEN),
    ]

    right_categories = [
        ("Platform & Virtualization", [
            ("Ubuntu Pro", "ESM, patching, compliance"),
            ("LXD / Incus", "Container management"),
            ("KVM / bhyve / VMM", "Hypervisor management"),
        ], GREEN),
        ("Notifications", [
            ("SMTP / Email", "Alert delivery"),
            ("Webhooks / Slack / Teams", "Real-time notifications"),
        ], DARK_BLUE),
        ("Firewalls", [
            ("UFW / iptables / nftables / PF", "Cross-platform firewall mgmt"),
        ], PRIMARY_BLUE),
    ]

    for col_offset, categories in [(0.0, left_categories), (6.3, right_categories)]:
        y = 1.35
        for cat_name, items, color in categories:
            cat_box = add_accent_rect(slide, 0.8 + col_offset, y, 2.6, 0.32, color)
            add_text_in_shape(cat_box, cat_name, size=9, color=WHITE, bold=True)

            for j, (name, desc) in enumerate(items):
                item_y = y + j * 0.32
                txBox = slide.shapes.add_textbox(
                    Inches(3.6 + col_offset), Inches(item_y), Inches(3.5), Inches(0.30),
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                set_run(tf.paragraphs[0], name, size=9, bold=True, color=DARK_BLUE)
                set_run(tf.paragraphs[0], f" \u2014 {desc}", size=8, bold=False, color=MEDIUM_GRAY)

            row_height = max(len(items) * 0.32, 0.36)
            y += row_height + 0.08

    # Bottom note — centered, larger font
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p,
            "SysManage integrates \u2014 not replaces \u2014 best-of-breed tools in each category",
            size=16, bold=True, color=DARK_TEXT)


def slide_18_why_security(prs):
    """Why Security Matters slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "Privileged access to every managed host \u2014 a compromised management platform is catastrophic",
        "Credential store \u2014 secrets engine holds sensitive keys and passwords",
        "Attack surface amplification \u2014 one server connects to hundreds of agents",
        "Compliance mandates \u2014 SOC2, HIPAA, FedRAMP require secure management tooling",
        "Supply chain risk \u2014 management software is a high-value target",
        "Audit requirements \u2014 every action must be traceable and tamper-evident",
        "Scalability testing in CI/CD \u2014 load testing in the pipeline to reduce DDoS attack impact",
    ]
    add_bullet_slide(slide, "Why Security Matters", bullets)


def slide_19_security_architecture(prs):
    """Security-First Architecture slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "mTLS agent communication \u2014 mutual TLS for all server-agent traffic",
        "JWT authentication with rotation \u2014 short-lived tokens, automatic refresh",
        "No inbound agent ports \u2014 agents connect outbound only, reducing attack surface",
        "UUID primary keys \u2014 prevents enumeration attacks on API endpoints",
        "OpenBAO encrypted secrets \u2014 HashiCorp Vault-compatible secrets management",
        "ECDSA P-521 license signing \u2014 cryptographically secure license validation",
        "TLS 1.2+ enforced \u2014 no legacy protocol support",
        "Cython-compiled Pro+ modules \u2014 source code protection for commercial features",
    ]
    add_bullet_slide(slide, "Security-First Architecture", bullets)


def slide_20_security_devprocess(prs):
    """Security in the Development Process slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_top_bar(slide)
    add_title_textbox(slide, "Security in the Development Process")
    add_bottom_bar(slide)
    add_footer_text(slide)

    tools = [
        ("Bandit", "Python security linter"),
        ("Semgrep Pro", "Static analysis rules"),
        ("SonarCloud", "Code quality & security"),
        ("CodeQL", "Semantic code analysis"),
        ("Safety", "Dependency vulnerability check"),
        ("Snyk", "Open source security scanning"),
        ("TruffleHog", "Secret detection in commits"),
        ("Black", "Deterministic code formatting"),
        ("PyLint 10/10", "Code quality enforcement"),
        ("Dependabot", "Automated dependency updates"),
        ("Pinned GH Actions", "SHA-pinned CI/CD actions"),
        ("Playwright", "End-to-end UI testing"),
    ]

    cols = 3
    for i, (tool, desc) in enumerate(tools):
        col = i % cols
        row = i // cols
        x = 0.8 + col * 4.0
        y = 1.5 + row * 1.2

        box = add_accent_rect(slide, x, y, 3.5, 0.9, PRIMARY_BLUE if i % 2 == 0 else DARK_BLUE)
        add_multiline_in_shape(box, [tool, desc], size=13, color=WHITE, alignment=PP_ALIGN.CENTER)


def slide_21_sbom(prs):
    """Software Bill of Materials (SBOM) slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_top_bar(slide)
    add_title_textbox(slide, "Software Bill of Materials (SBOM)")
    add_bottom_bar(slide)
    add_footer_text(slide)

    # Left column: what & why
    left_box = add_accent_rect(slide, 0.8, 1.5, 5.5, 3.5, DARK_BLUE)
    add_multiline_in_shape(left_box, [
        "What & Why",
        "",
        "\u2022  CycloneDX JSON format (v1.6)",
        "\u2022  Full dependency inventory for every release",
        "\u2022  Supply chain transparency & auditability",
        "\u2022  Separate SBOMs for backend (Python)",
        "    and frontend (Node.js) components",
        "\u2022  Agent packages include their own SBOM",
    ], size=13, color=WHITE, alignment=PP_ALIGN.LEFT)

    # Right column: how & where
    right_box = add_accent_rect(slide, 7.0, 1.5, 5.5, 3.5, PRIMARY_BLUE)
    add_multiline_in_shape(right_box, [
        "Generation & Distribution",
        "",
        "\u2022  Automated via CI/CD on every build",
        "\u2022  cyclonedx-bom (Python dependencies)",
        "\u2022  @cyclonedx/cyclonedx-npm (Node.js)",
        "\u2022  Published as GitHub release artifacts",
        "\u2022  Bundled inside every installer",
        "    (.deb, .rpm, .msi, .pkg, .snap, .tgz, ...)",
    ], size=13, color=WHITE, alignment=PP_ALIGN.LEFT)

    # Bottom: compatibility note
    compat_box = add_accent_rect(slide, 0.8, 5.4, 11.7, 1.0, BUTTON_GREEN)
    add_multiline_in_shape(compat_box, [
        "Compatible with: Grype  \u2022  Trivy  \u2022  Dependency-Track  \u2022  Snyk",
        "Inspect locally:  jq . backend-sbom.json  |  jq '.components | length'",
    ], size=13, color=WHITE, alignment=PP_ALIGN.CENTER)


def slide_22_demo_sequence(prs):
    """Demo Sequence slide — numbered walkthrough of the demo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_top_bar(slide)
    add_title_textbox(slide, "Demo Sequence")
    add_bottom_bar(slide)
    add_footer_text(slide)

    steps = [
        ("1", "Dashboard & Host Overview", "Navigate the main dashboard, view registered hosts and status"),
        ("2", "Host Detail Deep-Dive", "Live CPU/RAM/disk/network metrics, OS info, uptime, storage"),
        ("3", "Software Inventory & Updates", "View installed packages, detect available updates"),
        ("4", "Package Management", "Install and remove a package remotely from the browser"),
        ("5", "User Management & RBAC", "Create a user, assign roles (Admin/Operator/Viewer)"),
        ("6", "Settings & Configuration", "Server configuration, agent settings, preferences"),
        ("7", "Multi-Language Switching", "Switch UI language live across supported locales"),
    ]

    pro_steps = [
        ("8", "Health Analysis", "AI-powered system health scoring and recommendations"),
        ("9", "Compliance Audit", "Run CIS Benchmark / DISA STIG checks against a host"),
        ("10", "PDF Report Generation", "Generate and download a formatted compliance report"),
        ("11", "Secrets Management", "Store and retrieve secrets via OpenBAO integration"),
        ("12", "Container Management", "View and manage LXD containers on a host"),
        ("13", "CVE Vulnerability Scan", "Scan host packages against the CVE database"),
        ("14", "Alerting Configuration", "Configure Email, Webhook, Slack, and Teams alerts"),
        ("15", "Audit Log", "Browse tamper-evident audit trail with hash chain verification"),
    ]

    # Open Source section header
    oss_label = slide.shapes.add_textbox(Inches(0.8), Inches(1.25), Inches(3), Inches(0.3))
    tf = oss_label.text_frame
    set_run(tf.paragraphs[0], "Open Source", size=14, bold=True, color=PRIMARY_BLUE)

    # Pro+ section header
    pro_label = slide.shapes.add_textbox(Inches(6.5), Inches(1.25), Inches(3), Inches(0.3))
    tf2 = pro_label.text_frame
    set_run(tf2.paragraphs[0], "Pro+", size=14, bold=True, color=BUTTON_GREEN)

    for col_idx, items in enumerate([steps, pro_steps]):
        x = 0.8 + col_idx * 5.7
        for i, (num, title, desc) in enumerate(items):
            y = 1.55 + i * 0.68

            # Step number circle — larger for double-digit numbers
            circ_size = 0.45 if len(num) > 1 else 0.38
            circ = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x), Inches(y), Inches(circ_size), Inches(circ_size),
            )
            color = PRIMARY_BLUE if col_idx == 0 else BUTTON_GREEN
            circ.fill.solid()
            circ.fill.fore_color.rgb = color
            circ.line.fill.background()
            add_text_in_shape(circ, num, size=7, color=WHITE, bold=True)

            # Title + description
            txBox = slide.shapes.add_textbox(
                Inches(x + 0.55), Inches(y - 0.04), Inches(4.95), Inches(0.65),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            set_run(tf.paragraphs[0], title, size=13, bold=True, color=DARK_BLUE)
            p2 = tf.add_paragraph()
            p2.space_before = Pt(1)
            set_run(p2, desc, size=10, color=MEDIUM_GRAY)


def slide_23_live_demo(prs):
    """LIVE DEMO slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide)

    # Icon
    slide.shapes.add_picture(ICON_PNG, Inches(5.5), Inches(1.2), Inches(2.3))

    # LIVE DEMO text
    txBox = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, "LIVE DEMO", size=54, bold=True, color=WHITE)

    # Duration
    txBox2 = slide.shapes.add_textbox(Inches(3), Inches(5.2), Inches(7), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    set_run(p2, "30 Minutes", size=24, color=LIGHT_BLUE)

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.0), Inches(5.0), Inches(3.3), Inches(0.04),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_BLUE
    line.line.fill.background()


def slide_24_what_you_saw(prs):
    """What You Just Saw slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bullets = [
        "Real-time monitoring \u2014 live CPU, RAM, disk, network metrics via WebSocket",
        "Cross-platform agents \u2014 same management experience on Linux, macOS, Windows, *BSD",
        "One-click operations \u2014 package installs, user management, settings from the browser",
        "Pro+ intelligence layer \u2014 AI health analysis, CIS/STIG compliance, CVE scanning",
        "Enterprise-grade security \u2014 mTLS, tamper-evident audit logs, encrypted secrets",
        "Internationalized UI \u2014 seamless language switching across 14 languages",
    ]
    add_bullet_slide(slide, "What You Just Saw", bullets)


def slide_25_codebase_metrics(prs):
    """Codebase Metrics & Development History slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_top_bar(slide)
    add_title_textbox(slide, "Codebase & Development History")
    add_bottom_bar(slide)
    add_footer_text(slide)

    # Top-line stats row
    stats = [
        ("22 Months", "Under Development"),
        ("1M+ Lines", "of Code"),
        ("2,400+", "Commits"),
        ("72", "Release Tags"),
        ("4", "Repositories"),
    ]
    for i, (big, label) in enumerate(stats):
        x = 0.5 + i * 2.5
        box = add_accent_rect(slide, x, 1.3, 2.2, 0.9, DARK_BLUE)
        add_multiline_in_shape(box, [big, label], size=13, color=WHITE, bold=False,
                               alignment=PP_ALIGN.CENTER)

    # Repository timeline
    repos = [
        ("sysmanage", "Apr 2024", "Server: FastAPI + React", "1,047 commits"),
        ("sysmanage-agent", "Apr 2024", "Agent: Python cross-platform", "863 commits"),
        ("sysmanage-docs", "Sep 2025", "Docs site + packaging", "334 commits"),
        ("sysmanage-professional-plus", "Jan 2026", "Pro+ Cython modules", "171 commits"),
    ]

    headers = ["Repository", "Started", "Description", "Commits"]
    n_rows = len(repos) + 1
    table_shape = slide.shapes.add_table(n_rows, 4, Inches(0.5), Inches(2.5),
                                         Inches(12.0), Inches(0.35 * n_rows))
    table = table_shape.table
    col_ws = [3.5, 1.5, 5.0, 2.0]
    for ci, w in enumerate(col_ws):
        table.columns[ci].width = Inches(w)

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = FONT

    for r, (name, started, desc, commits) in enumerate(repos, start=1):
        for c, val in enumerate([name, started, desc, commits]):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else VERY_LIGHT_BLUE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = DARK_TEXT
                paragraph.font.name = FONT

    # Data model + test metrics row
    metrics = [
        ("Server DB", "74 tables \u2022 799 columns"),
        ("Agent DB", "7 tables"),
        ("Test Suite", "775 files \u2022 271K lines"),
    ]
    for i, (label, detail) in enumerate(metrics):
        x = 0.5 + i * 4.2
        box = add_accent_rect(slide, x, 4.8, 3.8, 0.7, PRIMARY_BLUE)
        add_multiline_in_shape(box, [label, detail], size=11, color=WHITE,
                               alignment=PP_ALIGN.CENTER)

    # Quarter-by-quarter velocity label
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.0), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    set_run(p, "Development velocity:  ", size=12, bold=True, color=DARK_BLUE)
    set_run(p, "Q2 2024 \u2192 project start  |  Q3-Q4 2025 \u2192 major feature build (730+ commits)  "
               "|  Q1 2026 \u2192 Pro+ launch + stabilization",
            size=11, bold=False, color=MEDIUM_GRAY)


def slide_26_roadmap_overview(prs):
    """Roadmap Overview slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_top_bar(slide)
    add_title_textbox(slide, "Roadmap Overview")
    add_bottom_bar(slide)
    add_footer_text(slide)

    # Historical releases (row 1) — compact
    history = [
        ("v0.9.0", "Foundation\n+ Core Platform"),
        ("v0.9.1", "CI/CD\n& Quality"),
        ("v0.9.2", "Mgmt\nFeatures"),
        ("v1.0.0", "Child Host\nFoundation"),
        ("v1.0.1", "Virtualization\nExpansion"),
        ("v1.0.2", "Platform\nMaturity"),
        ("v1.1.0", "Pro+\nLaunch"),
    ]

    for i, (ver, desc) in enumerate(history):
        x = 0.4 + i * 1.8
        box = add_accent_rect(slide, x, 1.25, 1.6, 1.1, MEDIUM_GRAY)
        add_multiline_in_shape(box, [ver, desc], size=8, color=WHITE, bold=False,
                               alignment=PP_ALIGN.CENTER)

    # "Shipped" label
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(1.05), Inches(3), Inches(0.2))
    tf = txBox.text_frame
    set_run(tf.paragraphs[0], "Shipped", size=8, bold=True, color=MEDIUM_GRAY)

    # Phases (rows 2-4)
    phases = [
        ("1", "Stabilization", "v1.2.0.0", "\u2713 Done", True),
        ("2", "Pro+ Prof.", "v1.3.0.0", "\u2713 Done", True),
        ("3", "AV + Firewall", "v1.4.0.0", "Apr 2026", False),
        ("4", "Stabilization", "v1.5.0.0", "May 2026", False),
        ("5", "Auto + Fleet", "v1.6.0.0", "May-Jun 2026", False),
        ("6", "Stabilization", "v1.7.0.0", "Jun 2026", False),
        ("7", "Stab. RC1", "v1.8.0.0", "Jun 2026", False),
        ("8", "Foundation", "v2.0.0.0", "Jul 2026", False),
        ("9", "Stab. RC2", "v2.1.0.0", "Aug 2026", False),
        ("10", "Virt+Obs+MFA", "v2.2.0.0", "Dec 2026", False),
        ("11", "Air-Gapped", "v2.3.0.0", "Jan 2027", False),
        ("12", "Federation", "v2.4.0.0", "Feb 2027", False),
        ("13", "Enterprise GA", "v3.0.0.0", "Mar 2027", False),
    ]

    for i, (num, desc, ver, status, done) in enumerate(phases):
        col = i % 7
        row = i // 7
        x = 0.3 + col * 1.82
        y = 2.7 + row * 1.85

        color = GREEN if done else PRIMARY_BLUE
        box = add_accent_rect(slide, x, y, 1.65, 1.55, color)
        add_multiline_in_shape(box, [
            f"Phase {num}",
            ver,
            "",
            desc,
            status,
        ], size=8, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

    # Target label
    txBox = slide.shapes.add_textbox(Inches(3.0), Inches(6.55), Inches(7), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, "Target: v3.0.0.0 Enterprise GA \u2014 Q1 2027", size=16, bold=True, color=DARK_BLUE)


def slide_27_futures_near(prs):
    """Futures: Near-Term slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    headers = ["Phase", "Version", "Description", "Target"]
    rows = [
        ("3", "v1.4.0.0", "AV Management + Firewall Orchestration", "Apr 2026"),
        ("4", "v1.5.0.0", "Stabilization", "May 2026"),
        ("5", "v1.6.0.0", "Automation + Fleet Management", "May-Jun 2026"),
    ]
    add_table_slide(slide, "Futures: Near-Term (Phases 3\u20135)", headers, rows,
                    col_widths=[1.0, 1.5, 6.0, 2.0])


def slide_28_futures_mid(prs):
    """Futures: Mid-Term slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    headers = ["Phase", "Version", "Description", "Target"]
    rows = [
        ("6", "v1.7.0.0", "Stabilization", "Jun 2026"),
        ("7", "v1.8.0.0", "Stabilization RC1", "Jun 2026"),
        ("8", "v2.0.0.0", "Foundation Features", "Jul 2026"),
        ("9", "v2.1.0.0", "Stabilization RC2", "Aug 2026"),
    ]
    add_table_slide(slide, "Futures: Mid-Term (Phases 6\u20139)", headers, rows,
                    col_widths=[1.0, 1.5, 6.0, 2.0])


def slide_29_futures_long(prs):
    """Futures: Long-Term slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    headers = ["Phase", "Version", "Description", "Target"]
    rows = [
        ("10", "v2.2.0.0", "Virtualization + Observability + MFA", "Dec 2026"),
        ("11", "v2.3.0.0", "Air-Gapped Support", "Jan 2027"),
        ("12", "v2.4.0.0", "Multi-Site Federation", "Feb 2027"),
        ("13", "v3.0.0.0", "Enterprise GA", "Mar 2027"),
    ]
    add_table_slide(slide, "Futures: Long-Term (Phases 10\u201313)", headers, rows,
                    col_widths=[1.0, 1.5, 6.0, 2.0])


def slide_30_takeaways(prs):
    """Key Takeaways slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_top_bar(slide)
    add_title_textbox(slide, "Key Takeaways")
    add_bottom_bar(slide)
    add_footer_text(slide)

    takeaways = [
        ("Cross-Platform", "6 operating systems \u2014 Linux, macOS, Windows, FreeBSD, OpenBSD, NetBSD"),
        ("Security-First", "mTLS, JWT rotation, tamper-evident auditing, encrypted secrets, zero inbound agent ports"),
        ("Open Source Core", "AGPLv3 \u2014 transparent, auditable, community-driven foundation"),
        ("Commercial Pro+", "Professional & Enterprise tiers for compliance, AI health, CVE scanning, and more"),
        ("Active Roadmap", "13-phase plan targeting v3.0.0.0 Enterprise GA in Q1 2027"),
    ]

    for i, (title, desc) in enumerate(takeaways):
        y = 1.5 + i * 1.05
        # Accent square
        sq = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y), Inches(0.15), Inches(0.7),
        )
        sq.fill.solid()
        sq.fill.fore_color.rgb = PRIMARY_BLUE
        sq.line.fill.background()

        # Title
        txBox = slide.shapes.add_textbox(Inches(1.3), Inches(y - 0.05), Inches(10), Inches(0.4))
        tf = txBox.text_frame
        set_run(tf.paragraphs[0], title, size=20, bold=True, color=DARK_BLUE)

        # Description
        txBox2 = slide.shapes.add_textbox(Inches(1.3), Inches(y + 0.35), Inches(10), Inches(0.4))
        tf2 = txBox2.text_frame
        set_run(tf2.paragraphs[0], desc, size=15, color=MEDIUM_GRAY)


def slide_31_partnerships(prs):
    """Partnerships slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_top_bar(slide)
    add_title_textbox(slide, "Partnership Opportunities")
    add_bottom_bar(slide)
    add_footer_text(slide)

    # Customers box
    cust_box = add_accent_rect(slide, 0.8, 1.5, 5.5, 3.5, DARK_BLUE)
    add_multiline_in_shape(cust_box, [
        "Customers",
        "",
        "\u2022  Organizations managing heterogeneous",
        "    OS estates (Linux, macOS, Windows, *BSD)",
        "\u2022  Compliance-driven environments",
        "    (SOC2, HIPAA, FedRAMP, CIS, DISA STIG)",
        "\u2022  Teams replacing expensive commercial",
        "    management tools with open source",
        "\u2022  MSPs and IT service providers",
    ], size=13, color=WHITE, alignment=PP_ALIGN.LEFT)

    # Channel partners box
    partner_box = add_accent_rect(slide, 7.0, 1.5, 5.5, 3.5, PRIMARY_BLUE)
    add_multiline_in_shape(partner_box, [
        "Channel Partners",
        "",
        "\u2022  Managed service providers (MSPs)",
        "\u2022  Systems integrators & consultancies",
        "\u2022  Value-added resellers (VARs)",
        "\u2022  Cloud hosting providers",
        "\u2022  Security & compliance consultancies",
        "\u2022  Regional technology distributors",
    ], size=13, color=WHITE, alignment=PP_ALIGN.LEFT)

    # Bottom note
    note_box = add_accent_rect(slide, 0.8, 5.5, 11.7, 1.0, BUTTON_GREEN)
    add_multiline_in_shape(note_box, [
        "We're actively seeking both direct customers and channel partners",
        "Open source core lowers barrier to entry  \u2022  Pro+ tiers create partner revenue opportunities",
    ], size=13, color=WHITE, alignment=PP_ALIGN.CENTER)


def slide_32_thank_you(prs):
    """Thank You slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide)

    # Logo
    slide.shapes.add_picture(LOGO_PNG, Inches(4.0), Inches(1.0), Inches(5.3))

    # Thank You
    txBox = slide.shapes.add_textbox(Inches(2), Inches(3.0), Inches(9), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, "Thank You", size=44, bold=True, color=WHITE)

    # URL
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    set_run(p2, "sysmanage.org", size=22, color=LIGHT_BLUE)

    # License
    txBox3 = slide.shapes.add_textbox(Inches(3), Inches(5.5), Inches(7), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    set_run(p3, "Licensed under AGPLv3", size=14, color=MEDIUM_GRAY)


# ── Main ───────────────────────────────────────────────────────────────

def main():
    print("Converting SVG logos to PNG...")
    convert_svgs()

    print("Creating presentation...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_builders = [
        slide_01_title,
        slide_02_agenda,
        slide_03_problem,
        slide_04_what_is,
        slide_05_who_bryan,
        slide_06_who_fedor,
        slide_07_growth_strategy,
        slide_08_architecture,
        slide_09_cross_platform,
        slide_10_packaging,
        slide_11_i18n,
        slide_12_oss_features,
        slide_13_oss_monitoring,
        slide_14_proplus_licensing,
        slide_15_proplus_professional,
        slide_16_proplus_enterprise,
        slide_17_ecosystem,
        slide_18_why_security,
        slide_19_security_architecture,
        slide_20_security_devprocess,
        slide_21_sbom,
        slide_22_demo_sequence,
        slide_23_live_demo,
        slide_24_what_you_saw,
        slide_25_codebase_metrics,
        slide_26_roadmap_overview,
        slide_27_futures_near,
        slide_28_futures_mid,
        slide_29_futures_long,
        slide_30_takeaways,
        slide_31_partnerships,
        slide_32_thank_you,
    ]

    for builder in slide_builders:
        name = builder.__name__
        print(f"  Building {name}...")
        builder(prs)

    print(f"Saving to {OUT_FILE}...")
    prs.save(OUT_FILE)
    print(f"Done! {len(prs.slides)} slides created.")
    print(f"Output: {OUT_FILE}")


if __name__ == "__main__":
    main()
